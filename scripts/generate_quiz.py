#!/usr/bin/env python3
"""
培训试题生成器脚本
根据培训材料生成标准格式的Excel试题文件
"""

import os
import sys
from typing import List, Dict, Any

# 尝试导入必要的库
try:
    import openpyxl
    from openpyxl import Workbook
except ImportError:
    print("请先安装 openpyxl: pip install openpyxl")
    sys.exit(1)


def create_quiz_excel(
    questions: List[Dict[str, Any]],
    output_path: str,
    category: str = "培训试题"
) -> str:
    """
    创建试题Excel文件
    
    支持单选题（选项A-D）、多选题（选项A-F，知识点超过4个时扩展E、F）、
    判断题（选项全部填"-"）。
    
    Args:
        questions: 试题列表，每道题是一个字典，支持 key:
                   type, question, option_a~f, answer, reference_answer, scoring_criteria
        output_path: 输出文件路径
        category: 试题分类名称
    
    Returns:
        输出文件的完整路径
    """
    wb = Workbook()
    ws = wb.active
    ws.title = "试题"
    
    # 写入表头（支持选项A-F，扩展多选题选项数量）
    headers = ['试题所属分类', '试题类型', '题目', '选项A', '选项B', '选项C', '选项D', '选项E', '选项F', '正确答案', '参考答案', '评分标准']
    ws.append(headers)
    
    # 写入试题
    for q in questions:
        row = [
            category,
            q.get('type', ''),
            q.get('question', ''),
            q.get('option_a', '-'),
            q.get('option_b', '-'),
            q.get('option_c', '-'),
            q.get('option_d', '-'),
            q.get('option_e', '-'),
            q.get('option_f', '-'),
            q.get('answer', ''),
            q.get('reference_answer', '-'),
            q.get('scoring_criteria', '-')
        ]
        ws.append(row)
    
    # 保存文件
    wb.save(output_path)
    return output_path


def generate_sample_questions(
    category: str = "培训主题",
    single_choice: int = 10,
    multiple_choice: int = 5,
    judgment: int = 5
) -> List[Dict[str, Any]]:
    """
    生成示例试题（用于测试）
    实际使用时应根据培训材料内容生成
    """
    questions = []
    
    # 单选题示例
    for i in range(single_choice):
        questions.append({
            'type': '单选题',
            'question': f'单选题示例 {i+1}：请根据培训内容选择正确答案',
            'option_a': '选项A内容',
            'option_b': '选项B内容（正确答案）',
            'option_c': '选项C内容',
            'option_d': '选项D内容',
            'answer': 'B'
        })
    
    # 多选题示例（展示4个选项和6个选项两种模式）
    # 4选项多选
    for i in range(min(2, multiple_choice)):
        questions.append({
            'type': '多选题',
            'question': f'多选题示例 {i+1}：以下哪些说法是正确的？',
            'option_a': '说法A（正确）',
            'option_b': '说法B（正确）',
            'option_c': '说法C（正确）',
            'option_d': '说法D',
            'option_e': '-',
            'option_f': '-',
            'answer': 'ABC'
        })
    
    # 6选项多选（知识点超过4个时使用）
    for i in range(2, multiple_choice):
        questions.append({
            'type': '多选题',
            'question': f'多选题示例 {i+1}（6选项）：以下哪些是值班管理的六大角色？',
            'option_a': '枢纽作用（Hub）',
            'option_b': '知识输出（Teacher）',
            'option_c': '目标负责（Target）',
            'option_d': '品质管控（Quality Guardian）',
            'option_e': '解决问题（Problem Solver）',
            'option_f': '内外沟通（Bridge）',
            'answer': 'ABCDEF'
        })
    
    # 判断题示例
    for i in range(judgment):
        is_correct = i % 2 == 0
        questions.append({
            'type': '判断题',
            'question': f'判断题示例 {i+1}：这是一个需要判断正误的陈述句。',
            'option_a': '-',
            'option_b': '-',
            'option_c': '-',
            'option_d': '-',
            'option_e': '-',
            'option_f': '-',
            'answer': '正确' if is_correct else '错误'
        })
    
    return questions


def read_pptx(file_path: str) -> str:
    """读取PPT文件内容"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        content = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = f"--- 幻灯片 {i} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_text += shape.text + "\n"
            content.append(slide_text)
        return "\n".join(content)
    except ImportError:
        print("请先安装 python-pptx: pip install python-pptx")
        return ""


def read_docx(file_path: str) -> str:
    """读取Word文件内容"""
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    except ImportError:
        print("请先安装 python-docx: pip install python-docx")
        return ""


def read_pdf(file_path: str) -> str:
    """读取PDF文件内容"""
    try:
        import pdfplumber
        content = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    content.append(text)
        return "\n".join(content)
    except ImportError:
        print("请先安装 pdfplumber: pip install pdfplumber")
        return ""


def read_training_material(file_path: str) -> str:
    """
    读取培训材料文件
    支持 .pptx, .docx, .pdf, .txt 等格式
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"文件不存在: {file_path}")
    
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pptx':
        return read_pptx(file_path)
    elif ext == '.docx':
        return read_docx(file_path)
    elif ext == '.pdf':
        return read_pdf(file_path)
    elif ext in ['.txt', '.md']:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


if __name__ == "__main__":
    # 示例：生成示例试题
    output_file = "示例试题.xlsx"
    questions = generate_sample_questions(
        category="示例培训",
        single_choice=10,
        multiple_choice=5,
        judgment=5
    )
    create_quiz_excel(questions, output_file, "示例培训")
    print(f"示例试题已生成: {output_file}")
