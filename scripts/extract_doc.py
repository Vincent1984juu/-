#!/usr/bin/env python3
"""
培训文档提取工具
支持 Word (.docx)、PDF (.pdf)、PPT (.pptx)
"""

import sys
import os
from pathlib import Path

def extract_docx(file_path):
    """提取Word文档内容"""
    try:
        from docx import Document
        doc = Document(file_path)
        content = []
        for para in doc.paragraphs:
            if para.text.strip():
                content.append(para.text.strip())
        # 提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    content.append(" | ".join(row_text))
        return "\n".join(content)
    except ImportError:
        return "[错误] 缺少 python-docx 库"
    except Exception as e:
        return f"[错误] 读取Word文档失败: {e}"

def extract_pdf(file_path):
    """提取PDF文档内容"""
    try:
        import pdfplumber
        content = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    content.append(text.strip())
        return "\n\n".join(content)
    except ImportError:
        # 使用pdftotext作为备用
        import subprocess
        try:
            result = subprocess.run(['pdftotext', '-layout', file_path, '-'], 
                                  capture_output=True, text=True, timeout=30)
            if result.returncode == 0:
                return result.stdout
            else:
                return f"[错误] pdftotext 失败: {result.stderr}"
        except Exception as e:
            return f"[错误] 读取PDF失败: {e}"
    except Exception as e:
        return f"[错误] 读取PDF失败: {e}"

def extract_pptx(file_path):
    """提取PPT文档内容"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            content.append(f"\n--- 幻灯片 {slide_num} ---\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text.strip())
        return "\n".join(content)
    except ImportError:
        return "[错误] 缺少 python-pptx 库"
    except Exception as e:
        return f"[错误] 读取PPT失败: {e}"

def extract_document(file_path):
    """根据文件类型自动选择提取方法"""
    file_path = Path(file_path)
    
    if not file_path.exists():
        return f"[错误] 文件不存在: {file_path}"
    
    suffix = file_path.suffix.lower()
    
    if suffix == '.docx':
        return extract_docx(file_path)
    elif suffix == '.pdf':
        return extract_pdf(file_path)
    elif suffix in ['.pptx', '.ppt']:
        return extract_pptx(file_path)
    else:
        return f"[错误] 不支持的文件格式: {suffix}"

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python3 extract_doc.py <文档路径>")
        print("支持格式: .docx, .pdf, .pptx")
        sys.exit(1)
    
    file_path = sys.argv[1]
    result = extract_document(file_path)
    print(result)
